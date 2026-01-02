import os
import io
import re
import math
import json
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
from fpdf import FPDF
from datetime import datetime
from typing import List, Dict, Any
import zlib
import time
import difflib
from duckduckgo_search import DDGS

from flask import Flask, render_template, request, redirect, url_for, flash, send_file, jsonify, session
from flask_sqlalchemy import SQLAlchemy
from flask_login import LoginManager, UserMixin, login_user, login_required, logout_user, current_user
from werkzeug.security import generate_password_hash, check_password_hash
from werkzeug.utils import secure_filename

from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity
import nltk
from nltk.tokenize import sent_tokenize
from humanizer import HumanizerEngine

# --- Initializations ---
app = Flask(__name__)
app.config['SECRET_KEY'] = 'plagicheck_secret_key_2026'
BASE_DIR = os.path.abspath(os.path.dirname(__file__))
app.config['SQLALCHEMY_DATABASE_URI'] = 'sqlite:///' + os.path.join(BASE_DIR, 'database.db')
app.config['UPLOAD_FOLDER'] = 'uploads'
app.config['REPORT_FOLDER'] = 'reports'

db = SQLAlchemy(app)
login_manager = LoginManager(app)
login_manager.login_view = 'login'

# Ensure folders exist
os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)
os.makedirs(app.config['REPORT_FOLDER'], exist_ok=True)

# NLTK initialization
try:
    nltk.data.find('tokenizers/punkt')
except LookupError:
    nltk.download('punkt')

# --- Database Models ---

class User(UserMixin, db.Model):
    id = db.Column(db.Integer, primary_key=True)
    username = db.Column(db.String(150), unique=True, nullable=False)
    password_hash = db.Column(db.String(200), nullable=False)
    full_name = db.Column(db.String(150))
    role = db.Column(db.String(20), default='user')  # 'admin' or 'user'
    status = db.Column(db.String(20), default='pending')  # 'pending', 'approved', 'rejected'

class DocumentRecord(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    user_id = db.Column(db.Integer, db.ForeignKey('user.id'), nullable=False)
    filename = db.Column(db.String(200), nullable=False)
    status = db.Column(db.String(50), default='Uploaded')  # 'Uploaded', 'Completed'
    raw_text = db.Column(db.Text)
    created_at = db.Column(db.DateTime, default=datetime.utcnow)
    
class ReportRecord(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    doc_id = db.Column(db.Integer, db.ForeignKey('document_record.id'), nullable=False)
    plag_score = db.Column(db.Float)
    ai_score = db.Column(db.Float)
    results_json = db.Column(db.Text)
    created_at = db.Column(db.DateTime, default=datetime.utcnow)

@login_manager.user_loader
def load_user(user_id):
    return User.query.get(int(user_id))

# --- Engines ---

class PlagiarismEngine:
    def __init__(self):
        self.user_agents = [
            "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36",
            "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) AppleWebKit/605.1.15 (KHTML, like Gecko) Version/17.2.1 Safari/605.1.15",
            "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36 Edg/120.0.0.0"
        ]

    def detect_language(self, text):
        # Heuristic: If > 10% of characters are Chinese, treat as Chinese
        if not text: return 'en'
        chinese_chars = sum(1 for char in text if '\u4e00' <= char <= '\u9fff')
        return 'zh' if (chinese_chars / len(text)) > 0.1 else 'en'

    def split_segments(self, text, lang):
        if lang == 'zh':
            # Split by Chinese punctuation
            return re.split(r'([。！？])', text)
        return sent_tokenize(text)

    def get_comparison_pool(self, current_doc_id, lang='en'):
        # Fetch all other completed docs
        # Implementation optimization: Only compare with same language to avoid bias
        others = DocumentRecord.query.filter(
            DocumentRecord.id != current_doc_id, 
            DocumentRecord.status == 'Completed'
        ).all()
        
        filtered = []
        for d in others:
            if self.detect_language(d.raw_text) == lang:
                filtered.append({"id": d.id, "text": d.raw_text, "filename": d.filename})
        return filtered

    def get_fingerprints(self, text, k=40):
        """Rolling K-Gram hashing for structural matching."""
        text = re.sub(r'\W+', '', text.lower())
        if len(text) < k: return set()
        return {zlib.adler32(text[i:i+k].encode()) for i in range(len(text) - k + 1)}

    def calculate_semantic_jaccard(self, text1, text2):
        """Precision keyword overlap to catch paraphrasing."""
        def get_keywords(t):
            return {w for w in re.findall(r'\w{4,}', t.lower()) if len(w) > 3}
        
        set1, set2 = get_keywords(text1), get_keywords(text2)
        if not set1 or not set2: return 0
        intersection = len(set1.intersection(set2))
        union = len(set1.union(set2))
        return intersection / union

    def search_web(self, query, lang='en'):
        """Perform a quick web search with professional browser identity."""
        try:
            import random
            ua = random.choice(self.user_agents)
            query_refined = query
            if lang == 'zh':
                query_refined += " lang:zh"
            
            with DDGS() as ddgs:
                results = list(ddgs.text(query_refined, max_results=8))
                if results:
                    bad_domains = [".cn", ".ru", ".tk", ".ml", ".ga", ".cf", "csdn", "zhihu", "blogspot", "wordpress"]
                    for r in results:
                        url = r.get("href", "").lower()
                        if any(bad in url for bad in bad_domains): continue
                        return {"title": r["title"], "url": r["href"]}
        except Exception as e:
            print(f"Web search error: {e}")
        return None

    def analyze(self, text, current_doc_id):
        lang = self.detect_language(text)
        segments = self.split_segments(text, lang)
        pool = self.get_comparison_pool(current_doc_id, lang)
        
        detailed_segments = []
        matches = {}
        total_plag_chars = 0
        
        # Pre-calculate Fingerprints for Pool
        pool_fingerprints = []
        for d in pool:
            pool_fingerprints.append({"id": d["id"], "fps": self.get_fingerprints(d["text"])})
        
        # --- TIER 1: Core Document-Wide Score (TF-IDF) ---
        core_index = 0
        vectorizer = None
        pool_vectors = None
        if pool:
            pool_texts = [d["text"] for d in pool]
            vectorizer = TfidfVectorizer().fit(pool_texts + [text])
            doc_vec = vectorizer.transform([text]).toarray()
            p_vectors = vectorizer.transform(pool_texts).toarray()
            sims = cosine_similarity(doc_vec, p_vectors)[0]
            core_index = round(float(sims.max()) * 100, 1)
            pool_vectors = p_vectors

        web_search_count = 0
        max_web_searches = 5 

        for seg in segments:
            clean_seg = seg.strip()
            if len(clean_seg) < 25:
                detailed_segments.append({"text": seg, "is_plag": False})
                continue
            
            is_matched = False
            best_sim = 0
            source_info = None

            # --- TIER 2: Fingerprint & Semantic Jaccard (Internal) ---
            seg_fps = self.get_fingerprints(clean_seg)
            if pool:
                for idx, p_item in enumerate(pool_fingerprints):
                    # 1. Digital Fingerprint Check
                    intersection = len(seg_fps.intersection(p_item["fps"]))
                    if intersection > 10: # High structural overlap
                        best_sim = 0.98
                        source = pool[idx]
                        source_info = {"id": source["id"], "title": source["filename"], "url": None}
                        is_matched = True
                        break
                    
                    # 2. Semantic Jaccard Backup (Catch Paraphrasing)
                    jaccard = self.calculate_semantic_jaccard(clean_seg, pool[idx]["text"])
                    if jaccard > 0.65:
                        best_sim = jaccard
                        source = pool[idx]
                        source_info = {"id": source["id"], "title": source["filename"], "url": None}
                        is_matched = True
                        break

            # --- TIER 3: Exact Match & Web ---
            if not is_matched and pool_vectors is not None:
                seg_vec = vectorizer.transform([seg]).toarray()
                sims = cosine_similarity(seg_vec, pool_vectors)[0]
                max_s = float(sims.max())
                if max_s > 0.85:
                    best_idx = int(sims.argmax())
                    source = pool[best_idx]
                    source_info = {"id": source["id"], "title": source["filename"], "url": None}
                    best_sim = max_s
                    is_matched = True

            # 3. Web Check (Phrase Scoping)
            if not is_matched and web_search_count < max_web_searches and len(clean_seg) > 85:
                phrase = " ".join(clean_seg.split()[:18])
                search_query = f'"{phrase}"'
                web_result = self.search_web(search_query, lang)
                if web_result:
                    source_info = {"id": f"web_{web_search_count}", "title": web_result["title"], "url": web_result["url"]}
                    best_sim = 0.99 
                    is_matched = True
                    web_search_count += 1
                    time.sleep(0.7)

            if is_matched:
                total_plag_chars += len(seg)
                detailed_segments.append({
                    "text": seg, "is_plag": True, "source_id": source_info["id"],
                    "similarity": round(best_sim * 100, 1), "url": source_info["url"]
                })
                m_key = source_info["url"] if source_info["url"] else source_info["id"]
                if m_key not in matches:
                    matches[m_key] = {"title": source_info["title"], "score": round(best_sim * 100, 1), "url": source_info["url"]}
                elif best_sim * 100 > matches[m_key]["score"]:
                    matches[m_key]["score"] = round(best_sim * 100, 1)
            else:
                detailed_segments.append({"text": seg, "is_plag": False})

        seg_plag_percentage = round((total_plag_chars / len(text)) * 100, 1) if text else 0
        final_score = round((core_index * 0.5) + (seg_plag_percentage * 0.5), 1)
        
        interpretation = "Original"
        if final_score > 70: interpretation = "Very High"
        elif final_score > 45: interpretation = "High"
        elif final_score > 25: interpretation = "Moderate"
        elif final_score > 10: interpretation = "Low"

        return {
            "score": final_score, "interpretation": interpretation,
            "segments": detailed_segments, "matches": list(matches.values())
        }
    

class AIEngine:
    def calculate_entropy(self, text):
        """Measures character-level predictability."""
        if not text: return 0
        probs = [text.count(c) / len(text) for c in set(text)]
        return -sum(p * math.log2(p) for p in probs)

    def calculate_compression(self, text):
        """Measures data complexity (AI text compresses very efficiently)."""
        if not text or len(text) < 30: return 0.85
        try:
            compressed = zlib.compress(text.encode('utf-8'))
            # Normalize by subtracting overhead
            return max(0.1, (len(compressed) - 11) / len(text))
        except:
            return 0.85

    def detect_language(self, text):
        if not text: return 'en'
        chinese_chars = sum(1 for char in text if '\u4e00' <= char <= '\u9fff')
        return 'zh' if (chinese_chars / len(text)) > 0.1 else 'en'

    def split_segments(self, text, lang):
        if lang == 'zh':
            parts = re.split(r'([。！？])', text)
            # Re-combine punctuation with the preceding segment
            segs = []
            for i in range(0, len(parts)-1, 2):
                segs.append(parts[i] + parts[i+1])
            if len(parts) % 2 == 1 and parts[-1]:
                segs.append(parts[-1])
            return segs
        return sent_tokenize(text)

    def analyze(self, text):
        lang = self.detect_language(text)
        segments = self.split_segments(text, lang)
        if not segments: return {"score": 0, "segments": []}
        
        # 1. Rolling Context & Document-Wide Analysis
        lengths = [len(s.split()) for s in segments if s.strip()]
        if not lengths: return {"score": 0, "segments": []}
        
        avg_len = sum(lengths) / len(lengths)
        
        # Rolling Burstiness Feature
        rolling_std = []
        if len(lengths) >= 3:
            for i in range(len(lengths) - 2):
                window = lengths[i:i+3]
                m = sum(window) / 3
                sd = math.sqrt(sum((x - m)**2 for x in window) / 3)
                rolling_std.append(sd)
        avg_rolling_std = sum(rolling_std) / len(rolling_std) if rolling_std else 0
        doc_burstiness = avg_rolling_std / (avg_len + 1e-6) # Professional threshold: AI < 0.25
        
        # 2. Advanced AI "Fingerprint" Dictionary (Expanded)
        ai_markers = [
            # Hedging & Cautious
            "it is important to note", "may suggest", "could be argued", "it is worth mentioning",
            "one must consider", "it can be argued", "research suggests", "it should be noted that",
            # Transitions & Explanatory
            "furthermore", "moreover", "consequently", "subsequently", "in addition",
            "as a result", "consequently", "thus", "notably", "significantly",
            "this indicates that", "it can be concluded", "it is evident that",
            "in order to fully understand", "as previously mentioned",
            # Summary & Conclusion
            "overall, it can be seen", "in conclusion", "to summarize", "in summary",
            "a key takeaway", "the following points", "this report aims to",
            # AI "Vibe" Words (Neutral/Substantial)
            "comprehensive", "essential", "pivotal", "nuanced", "delve", "leverage",
            "transformative", "underscore", "paradigm", "holistic", "strategic",
            "meticulous", "unwavering", "resilient", "vibrant", "dynamic", 
            "landscape", "pioneering", "ever-evolving", "realm"
        ]
        
        structural_patterns = [
            r"^(firstly|secondly|thirdly|finally|lastly|in addition|furthermore|moreover)[,:]",
            r"^\d+[\.\)]\s+", # Numbered lists
            r"^[•\-\*]\s+"   # Bullet points
        ]
        
        # Contractions check (AI rarely uses them in formal-adjacent text)
        contractions = [r"\b(can't|won't|don't|it's|i'm|he's|she's|they're|we're|you're)\b"]
        
        detailed = []
        total_ai_weighted_len = 0
        
        prev_entropy = None
        for i, seg in enumerate(segments):
            clean_seg = seg.strip()
            words = clean_seg.lower().split()
            if len(words) < 3:
                detailed.append({"text": seg, "is_ai": False, "prob": 0})
                continue
            
            # --- SIGNAL 1: Structural Uniformity ---
            dev = abs(len(words) - avg_len) / (avg_len + 5)
            uniformity = max(0, 1 - dev)
            
            # --- SIGNAL 2: Information Density (Entropy) ---
            entropy = self.calculate_entropy(clean_seg)
            
            # --- SIGNAL 2.1: Entropy Fluctuation (Human Variance) ---
            fluctuation_modifier = 0
            if prev_entropy is not None:
                diff = abs(entropy - prev_entropy)
                # Humans fluctuate entropy between sentences; AI is remarkably stable.
                if diff < 0.15: # Too stable
                    fluctuation_modifier = 0.15
                elif diff > 0.5: # Human-like variance
                    fluctuation_modifier = -0.1
            prev_entropy = entropy

            if lang == 'zh':
                if 4.8 <= entropy <= 5.8:
                    entropy_score = 0.85
                elif entropy < 4.8:
                    entropy_score = 0.6
                else: entropy_score = max(0, 1 - (entropy - 5.7))
            else:
                if 3.5 <= entropy <= 4.3:
                    entropy_score = 0.85
                elif entropy < 3.5:
                    entropy_score = 0.6
                else: entropy_score = max(0, 1 - (entropy - 4.2))
            
            # --- SIGNAL 3: Pattern Repeating (Compression) ---
            comp_ratio = self.calculate_compression(clean_seg)
            if lang == 'zh':
                comp_score = max(0, min(1, (0.85 - comp_ratio) / 0.25))
            else:
                comp_score = max(0, min(1, (0.75 - comp_ratio) / 0.35))
            
            # --- SIGNAL 4: Vocabulary Markers & Connective Density ---
            marker_count = 0
            connective_density = 0
            if lang == 'en':
                for marker in ai_markers:
                    if marker in clean_seg.lower():
                        marker_count += 1
                # Connectives are structural glue overused by AI
                connectives = ["therefore", "however", "additionally", "furthermore", "consequently"]
                connective_count = sum(1 for c in connectives if c in words)
                connective_density = min(connective_count * 0.25, 0.5)

            marker_score = min(marker_count * 0.2, 0.65) if lang == 'en' else 0
            
            # --- SIGNAL 5: List/Report Structure ---
            struct_score = 0
            if lang == 'en':
                for pattern in structural_patterns:
                    if re.search(pattern, clean_seg, re.IGNORECASE):
                        struct_score = 0.4
                        break
            elif lang == 'zh':
                if re.match(r'^[\d一二三四五].*[：:]', clean_seg):
                    struct_score = 0.3

            # --- SIGNAL 6: Formal/Contraction Check ---
            casual_score = 0
            if lang == 'en':
                for pattern in contractions:
                    if re.search(pattern, clean_seg, re.IGNORECASE):
                        casual_score = -0.2
                        break
            
            # --- SIGNAL 7: Repetitive Explanation ---
            redundancy_score = 0
            if lang == 'en' and len(words) > 15:
                unique_words = set(words)
                if len(unique_words) / len(words) < 0.6:
                    redundancy_score = 0.2
            elif lang == 'zh' and len(clean_seg) > 20:
                unique_chars = set(clean_seg)
                if len(unique_chars) / len(clean_seg) < 0.45:
                    redundancy_score = 0.2
            
            # --- ULTIMATE ENSEMBLE SCORING ---
            # AI burstiness is typically low (<0.25 normalized)
            burst_modifier = 0.2 if doc_burstiness < 0.25 else -0.1
            
            prob_raw = (uniformity * 0.1) + \
                       (entropy_score * 0.25) + \
                       (comp_score * 0.15) + \
                       (burst_modifier) + \
                       (fluctuation_modifier) + \
                        connective_density + \
                       marker_score + \
                       struct_score + \
                       casual_score + \
                       redundancy_score
            
            prob = round(min(max(prob_raw, 0.01), 0.99) * 100, 1)
            is_ai = prob > 48.0 # Refined threshold with more high-signal inputs
            
            if is_ai:
                total_ai_weighted_len += (prob/100.0) * len(seg)
            
            detailed.append({"text": seg, "is_ai": is_ai, "prob": prob})
            
        doc_score = round((total_ai_weighted_len / len(text)) * 100, 1) if text else 0
        if len(text) < 200:
            doc_score = round(doc_score * 0.7, 1)
            
        return {"score": doc_score, "segments": detailed}

# --- Engine Instances ---
plag_engine = PlagiarismEngine()
ai_engine = AIEngine()
humanizer_engine = HumanizerEngine()

# --- Utils ---

def extract_text(filepath):
    ext = filepath.split('.')[-1].lower()
    text = ""
    try:
        if ext == 'pdf':
            doc = fitz.open(filepath)
            for page in doc: text += page.get_text()
            doc.close()
        elif ext in ['docx', 'doc']:
            doc = Document(filepath)
            text = "\n".join([p.text for p in doc.paragraphs])
        elif ext in ['pptx', 'ppt']:
            prs = Presentation(filepath)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"): text += shape.text + "\n"
        else:
            with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                text = f.read()
    except Exception as e:
        print(f"Extraction error: {e}")
    return text

import unicodedata
import re

def clean_text_for_pdf(text):
    if not isinstance(text, str):
        return ""

    # Normalize Unicode (smart quotes, accents, etc.)
    text = unicodedata.normalize("NFKD", text)

    # Replace problematic characters
    replacements = {
        "–": "-",
        "—": "-",
        "―": "-",
        "“": '"',
        "”": '"',
        "„": '"',
        "‘": "'",
        "’": "'",
        "‚": "'",
        "\u00A0": " ",   # non-breaking space
        "\u200b": "",    # zero-width space
        "\u2028": " ",
        "\u2029": " ",
    }

    for bad, good in replacements.items():
        text = text.replace(bad, good)

    # Remove emojis & non-BMP characters
    text = re.sub(r"[^\u0000-\uFFFF]", "", text)

    # Remove control characters
    text = re.sub(r"[\x00-\x08\x0B\x0C\x0E-\x1F]", "", text)

    return text.strip()

class ModernPDF(FPDF):
    def footer(self):
        self.set_y(-15)
        self.set_font("DejaVu", "", 8)
        self.set_text_color(150, 150, 150)
        self.cell(0, 10, f"Page {self.page_no()} | Generated by PlagiCheck AI Labs", 0, 0, "C")

def generate_pdf_report(report_id, username, report_type='similarity'):
    report = ReportRecord.query.get(report_id)
    doc = DocumentRecord.query.get(report.doc_id)
    results = json.loads(report.results_json)
    
    # Get interpretation if available
    interpretation = results.get("plagiarism", {}).get("interpretation", "Original") if report_type == 'similarity' else "N/A"

    pdf = ModernPDF()
    pdf.set_margins(15, 15, 15)
    pdf.add_page()
    pdf.add_font("DejaVu", "", "font1/DejaVuSans.ttf", uni=True)
    pdf.add_font("DejaVu", "B", "font1/DejaVuSans.ttf", uni=True) # Faking bold if necessary or using same file

    # ===== MODERN HEADER BAR =====
    if report_type == 'ai':
        primary_color = (37, 99, 235) # AI Blue
        title = "AI CONTENT ANALYSIS REPORT"
    else:
        primary_color = (220, 38, 38) # Plag Red
        title = "SIMILARITY ANALYSIS REPORT"

    pdf.set_fill_color(*primary_color)
    pdf.rect(0, 0, 210, 40, "F")
    
    pdf.set_xy(10, 12)
    pdf.set_font("DejaVu", "B", 22)
    pdf.set_text_color(255, 255, 255)
    pdf.cell(0, 10, "PlagiCheck AI", ln=True)
    
    pdf.set_font("DejaVu", "", 9)
    pdf.set_xy(10, 22)
    pdf.cell(0, 10, title, ln=True)

    pdf.ln(20)

    # ===== METADATA GRID =====
    pdf.set_y(50)
    pdf.set_text_color(50, 50, 50)
    pdf.set_font("DejaVu", "B", 10)
    
    # Column 1
    pdf.set_x(15)
    pdf.cell(30, 8, "DOCUMENT:", ln=0)
    pdf.set_font("DejaVu", "", 10)
    pdf.cell(100, 8, clean_text_for_pdf(doc.filename), ln=1)
    
    # Column 2
    pdf.set_font("DejaVu", "B", 10)
    pdf.set_x(15)
    pdf.cell(30, 8, "AUTHOR:", ln=0)
    pdf.set_font("DejaVu", "", 10)
    pdf.cell(60, 8, clean_text_for_pdf(username), ln=0)
    
    pdf.set_font("DejaVu", "B", 10)
    pdf.cell(20, 8, "DATE:", ln=0)
    pdf.set_font("DejaVu", "", 10)
    pdf.cell(40, 8, doc.created_at.strftime('%Y-%m-%d'), ln=1)

    pdf.ln(10)
    pdf.line(15, pdf.get_y(), 195, pdf.get_y())
    pdf.ln(10)

    # ===== SCORE BADGE SECTION =====
    pdf.set_fill_color(252, 252, 253)
    pdf.rect(15, pdf.get_y(), 180, 45, "F")
    pdf.set_draw_color(226, 232, 240)
    pdf.rect(15, pdf.get_y(), 180, 45, "D")
    
    start_y = pdf.get_y()
    
    # Large Score
    pdf.set_xy(25, start_y + 10)
    pdf.set_font("DejaVu", "B", 34)
    pdf.set_text_color(*primary_color)
    score = report.ai_score if report_type == 'ai' else report.plag_score
    pdf.cell(60, 25, f"{score}%", ln=0, align="L")
    
    # Label & Interpretation
    pdf.set_xy(85, start_y + 12)
    pdf.set_font("DejaVu", "B", 12)
    pdf.set_text_color(71, 85, 105)
    main_label = "AI PROBABILITY" if report_type == 'ai' else "SIMILARITY INDEX"
    pdf.cell(0, 8, main_label, ln=1)
    
    pdf.set_x(85)
    pdf.set_font("DejaVu", "", 10)
    if report_type == 'similarity':
        status_text = f"Status: {interpretation} Matching"
    else:
        status_text = "Linguistic pattern analysis complete."
    pdf.cell(0, 8, status_text, ln=1)

    pdf.set_y(start_y + 55)

    # ===== ANALYSIS HIGHLIGHTS =====
    pdf.set_font("DejaVu", "B", 14)
    pdf.set_text_color(30, 41, 59)
    pdf.cell(0, 10, "Detailed Analysis Breakdown", ln=True)
    pdf.ln(5)

    pdf.set_font("DejaVu", "", 10)
    pdf.set_text_color(0, 0, 0)

    if report_type == 'ai':
        segments = results.get("ai", {}).get("segments", [])
        for seg in segments:
            safe_text = clean_text_for_pdf(seg.get("text", ""))
            if not safe_text: continue
            
            if seg.get("is_ai"):
                pdf.set_fill_color(239, 246, 255)
                pdf.set_text_color(*primary_color)
                # Subtle indicator
                pdf.set_font("DejaVu", "B", 8)
                pdf.cell(180, 5, f"[AI Confidence: {seg.get('prob')}%]", ln=True)
                pdf.set_font("DejaVu", "", 10)
                pdf.multi_cell(180, 6, safe_text, fill=True)
            else:
                pdf.set_text_color(0, 0, 0)
                pdf.multi_cell(180, 6, safe_text)
            pdf.ln(2)
    else:
        segments = results.get("plagiarism", {}).get("segments", [])
        for seg in segments:
            safe_text = clean_text_for_pdf(seg.get("text", ""))
            if not safe_text: continue

            if seg.get("is_plag"):
                pdf.set_fill_color(254, 242, 242)
                pdf.set_text_color(*primary_color)
                pdf.set_font("DejaVu", "B", 8)
                match_header = f"[EXACT MATCH: {seg.get('similarity')}%]"
                if seg.get("url"):
                    match_header += f" | Source: {seg.get('url')}"
                pdf.multi_cell(180, 5, match_header, fill=True)
                
                pdf.set_font("DejaVu", "", 10)
                pdf.multi_cell(180, 6, safe_text, fill=True)
            else:
                pdf.set_text_color(0, 0, 0)
                pdf.multi_cell(180, 6, safe_text)
            pdf.ln(2)

    # ===== SAVE FILE =====
    suffix = f"{report_type}_{report_id}"
    filepath = os.path.join(
        app.config["REPORT_FOLDER"],
        f"report_{suffix}.pdf"
    )

    pdf.output(filepath)
    return filepath

# --- Routes ---

@app.route('/')
def index():
    if current_user.is_authenticated:
        if current_user.role == 'admin': return redirect(url_for('admin_dashboard'))
        return redirect(url_for('user_dashboard'))
    return redirect(url_for('login'))

@app.route('/login', methods=['GET', 'POST'])
def login():
    if request.method == 'POST':
        username = request.form.get('username')
        password = request.form.get('password')
        user = User.query.filter_by(username=username).first()
        if user and check_password_hash(user.password_hash, password):
            if user.status != 'approved' and user.role != 'admin':
                flash('Your account is pending admin approval.', 'warning')
                return redirect(url_for('login'))
            login_user(user)
            return redirect(url_for('index'))
        flash('Invalid credentials or unapproved account.', 'danger')
    return render_template('login.html')

@app.route('/register', methods=['GET', 'POST'])
def register():
    if request.method == 'POST':
        username = request.form.get('username')
        password = request.form.get('password')
        fullname = request.form.get('fullname')
        if User.query.filter_by(username=username).first():
            flash('Username already exists.', 'danger')
        else:
            new_user = User(username=username, password_hash=generate_password_hash(password), full_name=fullname)
            db.session.add(new_user)
            db.session.commit()
            flash('Registration successful! Please wait for admin approval.', 'success')
            return redirect(url_for('login'))
    return render_template('register.html')

@app.route('/logout')
@login_required
def logout():
    logout_user()
    return redirect(url_for('login'))

@app.route('/dashboard')
@login_required
def user_dashboard():
    docs = DocumentRecord.query.filter_by(user_id=current_user.id).order_by(DocumentRecord.created_at.desc()).all()
    reports = {r.doc_id: r for r in ReportRecord.query.all()}
    return render_template('user_dashboard.html', docs=docs, reports=reports)

@app.route('/admin')
@login_required
def admin_dashboard():
    if current_user.role != 'admin': return redirect(url_for('index'))
    pending_users = User.query.filter_by(status='pending', role='user').all()
    all_users = User.query.filter_by(role='user').all()
    all_docs = DocumentRecord.query.all()
    reports = {r.doc_id: r for r in ReportRecord.query.all()}
    return render_template('admin_dashboard.html', pending=pending_users, users=all_users, docs=all_docs, reports=reports)

@app.route('/admin/approve/<int:user_id>')
@login_required
def approve_user(user_id):
    if current_user.role != 'admin': return redirect(url_for('index'))
    user = User.query.get(user_id)
    if user:
        user.status = 'approved'
        db.session.commit()
        flash(f'User {user.username} approved.', 'success')
    return redirect(url_for('admin_dashboard'))

@app.route('/admin/reject/<int:user_id>')
@login_required
def reject_user(user_id):
    if current_user.role != 'admin': return redirect(url_for('index'))
    user = User.query.get(user_id)
    if user:
        user.status = 'rejected'
        db.session.commit()
        flash(f'User {user.username} rejected.', 'info')
    return redirect(url_for('admin_dashboard'))

@app.route('/admin/delete_user/<int:user_id>')
@login_required
def delete_user(user_id):
    if current_user.role != 'admin': return redirect(url_for('index'))
    user = User.query.get(user_id)
    if user:
        db.session.delete(user)
        db.session.commit()
        flash('User deleted.', 'success')
    return redirect(url_for('admin_dashboard'))

@app.route('/upload', methods=['POST'])
@login_required
def upload():
    # Handle text input
    text_content = request.form.get('text_content')
    if text_content and text_content.strip():
        # Save text as a temporary .txt file
        filename = f"text_input_{datetime.now().strftime('%Y%m%d_%H%M%S')}.txt"
        filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        with open(filepath, 'w', encoding='utf-8') as f:
            f.write(text_content)
        
        new_doc = DocumentRecord(user_id=current_user.id, filename=filename, raw_text=text_content)
        db.session.add(new_doc)
        db.session.commit()
        
        flash('Text uploaded successfully. You can now run the analysis.', 'success')
        return redirect(url_for('user_dashboard'))
    
    # Handle file upload
    file = request.files.get('file')
    if file and file.filename:
        filename = secure_filename(file.filename)
        filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        file.save(filepath)
        
        text = extract_text(filepath)
        new_doc = DocumentRecord(user_id=current_user.id, filename=filename, raw_text=text)
        db.session.add(new_doc)
        db.session.commit()
        
        flash('Document uploaded. You can now run the analysis.', 'success')
    return redirect(url_for('user_dashboard'))

@app.route('/run/<int:doc_id>')
@login_required
def run_scan(doc_id):
    doc = DocumentRecord.query.get(doc_id)
    if not doc or (doc.user_id != current_user.id and current_user.role != 'admin'):
        flash('Access denied.', 'danger')
        return redirect(url_for('index'))
    
    plag_eng = PlagiarismEngine()
    ai_eng = AIEngine()
    
    plag_res = plag_eng.analyze(doc.raw_text, doc.id)
    ai_res = ai_eng.analyze(doc.raw_text)
    
    results = {"plagiarism": plag_res, "ai": ai_res}
    
    existing_report = ReportRecord.query.filter_by(doc_id=doc.id).first()
    if existing_report:
        existing_report.plag_score = plag_res["score"]
        existing_report.ai_score = ai_res["score"]
        existing_report.results_json = json.dumps(results)
    else:
        report = ReportRecord(doc_id=doc.id, plag_score=plag_res["score"], ai_score=ai_res["score"], results_json=json.dumps(results))
        db.session.add(report)
    
    doc.status = 'Completed'
    db.session.commit()
    flash('Analysis complete.', 'success')
    return redirect(url_for('user_dashboard'))

@app.route('/report/<int:doc_id>')
@login_required
def view_report(doc_id):
    doc = DocumentRecord.query.get_or_404(doc_id)
    report = ReportRecord.query.filter_by(doc_id=doc_id).first()
    if not report: 
        flash('Report not found.', 'danger')
        return redirect(url_for('user_dashboard'))
    
    results = json.loads(report.results_json)
    
    # Combine segments for easier template rendering
    plag_segs = results.get('plagiarism', {}).get('segments', [])
    ai_segs = results.get('ai', {}).get('segments', [])
    
    combined = []
    for i in range(len(plag_segs)):
        p = plag_segs[i]
        a = ai_segs[i] if i < len(ai_segs) else {}
        combined.append({
            'text': p.get('text', ''),
            'is_plag': p.get('is_plag', False),
            'similarity': p.get('similarity', 0),
            'url': p.get('url', ''),
            'is_ai': a.get('is_ai', False),
            'prob': a.get('prob', 0)
        })
    
    matches = results.get('plagiarism', {}).get('matches', [])
    interpretation = results.get('plagiarism', {}).get('interpretation', 'Original')
    
    return render_template('report.html', 
                           doc=doc, 
                           report=report, 
                           combined_segments=combined,
                           matches=matches,
                           interpretation=interpretation)

@app.route('/download/<int:report_id>')
@login_required
def download_report(report_id):
    path = generate_pdf_report(report_id, current_user.username, 'similarity')
    return send_file(path, as_attachment=True)

@app.route('/download_similarity/<int:report_id>')
@login_required
def download_similarity(report_id):
    path = generate_pdf_report(report_id, current_user.username, 'similarity')
    return send_file(path, as_attachment=True)

@app.route('/download_ai/<int:report_id>')
@login_required
def download_ai(report_id):
    path = generate_pdf_report(report_id, current_user.username, 'ai')
    return send_file(path, as_attachment=True)

@app.route('/delete_doc/<int:doc_id>')
@login_required
def delete_doc(doc_id):
    doc = DocumentRecord.query.get(doc_id)
    if doc and (doc.user_id == current_user.id or current_user.role == 'admin'):
        ReportRecord.query.filter_by(doc_id=doc_id).delete()
        db.session.delete(doc)
        db.session.commit()
        flash('Document and report deleted.', 'info')
    return redirect(url_for('user_dashboard'))

# --- Main ---

if __name__ == '__main__':
    with app.app_context():
        db.create_all()
        # Bootstrap Admin
        if not User.query.filter_by(username='LaiTech001').first():
            admin = User(
                username='LaiTech001',
                password_hash=generate_password_hash('LaiTech001*'),
                full_name='System Administrator',
                role='admin',
                status='approved'
            )
            db.session.add(admin)
            db.session.commit()
            print("Admin user LaiTech001 created.")
    app.run(host='0.0.0.0', port=5000, debug=True)
